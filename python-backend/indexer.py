# -*- coding: utf-8 -*-
"""
파일 인덱싱 엔진 - 비동기 파일 시스템 크롤링 및 텍스트 추출
"""

import os
import threading
from pathlib import Path
from typing import List, Callable, Optional, Dict
import logging
import time
from queue import Queue
from datetime import datetime
import traceback
import signal
from functools import wraps
import re
import unicodedata
import shutil
import tempfile

# 텍스트 추출 라이브러리
import chardet  # 인코딩 자동 감지

# 사용자 입력 감지 (키보드/마우스)
try:
    from pynput import mouse, keyboard
    PYNPUT_AVAILABLE = True
except ImportError:
    PYNPUT_AVAILABLE = False
    logging.warning("pynput not installed. User activity monitoring disabled.")

# 문서 파일 파싱
try:
    import docx  # python-docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not installed. .docx support disabled.")

try:
    from pptx import Presentation  # python-pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not installed. .pptx support disabled.")

try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logging.warning("PyMuPDF not installed. .pdf support disabled.")

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logging.warning("openpyxl not installed. .xlsx support disabled.")

try:
    import win32com.client
    import pythoncom
    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False
    logging.warning("pywin32 not installed. .doc, .ppt, .xls, .hwp support disabled.")

try:
    import olefile
    OLEFILE_AVAILABLE = True
except ImportError:
    OLEFILE_AVAILABLE = False
    logging.warning("olefile not installed. Alternative .hwp support disabled.")

from database import DatabaseManager
import sys
import io

# ========================================
# UTF-8 전역 설정 (최우선 실행)
# ========================================
# Windows 콘솔 코드 페이지를 UTF-8로 설정
if sys.platform == 'win32':
    try:
        import ctypes
        kernel32 = ctypes.windll.kernel32
        kernel32.SetConsoleCP(65001)
        kernel32.SetConsoleOutputCP(65001)
    except Exception:
        pass

# stdout/stderr를 UTF-8로 재설정 (안전하게)
try:
    if hasattr(sys.stdout, 'buffer') and sys.stdout.encoding != 'utf-8':
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace', line_buffering=True)
except Exception:
    pass

try:
    if hasattr(sys.stderr, 'buffer') and sys.stderr.encoding != 'utf-8':
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace', line_buffering=True)
except Exception:
    pass

# 로그 디렉토리 생성
LOG_DIR = os.path.join(os.path.dirname(__file__), 'logs')
os.makedirs(LOG_DIR, exist_ok=True)

# 로깅 설정 (콘솔 + 파일, UTF-8 인코딩 강제)
log_file = os.path.join(LOG_DIR, 'indexer.log')
logging.basicConfig(
    level=logging.INFO, 
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(sys.stdout),
        logging.FileHandler(log_file, encoding='utf-8', mode='a')
    ]
)
# UTF-8 인코딩 재확인
for handler in logging.root.handlers:
    if isinstance(handler, logging.StreamHandler) and hasattr(handler.stream, 'reconfigure'):
        try:
            handler.stream.reconfigure(encoding='utf-8', errors='replace')
        except Exception:
            pass
        
logger = logging.getLogger(__name__)

# 상수 정의
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
PARSE_TIMEOUT = 60  # 60초


class TimeoutError(Exception):
    """타임아웃 예외"""
    pass


def timeout_handler(signum, frame):
    """타임아웃 시그널 핸들러"""
    raise TimeoutError("Parsing timeout")


def with_timeout(seconds):
    """타임아웃 데코레이터 (Windows에서는 제한적)"""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            # Windows에서는 signal.alarm이 작동하지 않으므로 threading 사용
            result = [TimeoutError("Timeout")]
            
            def target():
                try:
                    result[0] = func(*args, **kwargs)
                except Exception as e:
                    result[0] = e
            
            thread = threading.Thread(target=target, daemon=True)
            thread.start()
            thread.join(timeout=seconds)
            
            if thread.is_alive():
                # 타임아웃 발생
                raise TimeoutError(f"Function execution exceeded {seconds} seconds")
            
            if isinstance(result[0], Exception):
                raise result[0]
            
            return result[0]
        return wrapper
    return decorator


class UserActivityMonitor:
    """
    사용자 활동 모니터 (키보드/마우스)
    
    Windows 전체 시스템에서 키보드 입력 및 마우스 움직임을 감지하고,
    활동이 감지되면 인덱싱을 일시정지합니다.
    """
    
    def __init__(self, idle_threshold: float = 2.0):
        """
        Args:
            idle_threshold: 사용자 활동이 없는 시간 (초) - 이 시간이 지나면 인덱싱 재개
        """
        self.idle_threshold = idle_threshold
        self.last_activity_time = time.time()
        self.is_monitoring = False
        self.keyboard_listener = None
        self.mouse_listener = None
        self._lock = threading.Lock()
        
    def start(self):
        """모니터링 시작"""
        if not PYNPUT_AVAILABLE:
            logger.warning("pynput이 설치되지 않아 사용자 활동 모니터링을 사용할 수 없습니다.")
            return False
        
        if self.is_monitoring:
            return True
        
        try:
            # 키보드 리스너
            self.keyboard_listener = keyboard.Listener(
                on_press=self._on_activity,
                on_release=self._on_activity
            )
            
            # 마우스 리스너
            self.mouse_listener = mouse.Listener(
                on_move=self._on_activity,
                on_click=self._on_activity,
                on_scroll=self._on_activity
            )
            
            self.keyboard_listener.start()
            self.mouse_listener.start()
            self.is_monitoring = True
            self.last_activity_time = time.time()
            
            logger.info(f"사용자 활동 모니터링 시작 (대기 시간: {self.idle_threshold}초)")
            return True
        
        except Exception as e:
            logger.error(f"사용자 활동 모니터링 시작 실패: {e}")
            return False
    
    def stop(self):
        """모니터링 중지"""
        if not self.is_monitoring:
            return
        
        try:
            if self.keyboard_listener:
                self.keyboard_listener.stop()
            if self.mouse_listener:
                self.mouse_listener.stop()
            
            self.is_monitoring = False
            logger.info("사용자 활동 모니터링 중지됨")
        
        except Exception as e:
            logger.error(f"사용자 활동 모니터링 중지 오류: {e}")
    
    def _on_activity(self, *args, **kwargs):
        """사용자 활동 감지 시 호출"""
        with self._lock:
            self.last_activity_time = time.time()
    
    def is_user_active(self) -> bool:
        """
        사용자가 활동 중인지 확인
        
        Returns:
            True면 최근 활동 있음 (대기 필요), False면 유휴 상태 (인덱싱 진행 가능)
        """
        with self._lock:
            elapsed = time.time() - self.last_activity_time
            return elapsed < self.idle_threshold
    
    def wait_until_idle(self, check_interval: float = 0.05, stop_flag: threading.Event = None) -> bool:
        """
        사용자가 유휴 상태가 될 때까지 대기
        
        Args:
            check_interval: 체크 간격 (초) - 기본 0.05초로 매우 빠른 반응
            stop_flag: 중지 플래그 (인덱싱 중단 시 사용)
        
        Returns:
            True면 유휴 상태 도달, False면 중지 요청됨
        """
        while self.is_user_active():
            if stop_flag and stop_flag.is_set():
                return False
            
            # 사용자 활동 중 - 대기 (매우 짧은 간격으로 체크)
            time.sleep(check_interval)
        
        return True
    
    def get_idle_time(self) -> float:
        """마지막 활동 이후 경과 시간 (초) 반환"""
        with self._lock:
            return time.time() - self.last_activity_time


class FileIndexer:
    """파일 인덱싱 엔진 - Worker Thread에서 실행"""
    
    # 지원하는 파일 확장자
    SUPPORTED_TEXT_EXTENSIONS = {'.txt', '.log', '.md', '.py', '.js', '.ts', '.jsx', '.tsx', 
                                  '.java', '.cpp', '.c', '.h', '.cs', '.json', '.xml', '.html', 
                                  '.css', '.sql', '.sh', '.bat', '.ps1', '.yaml', '.yml'}
    
    SUPPORTED_DOC_EXTENSIONS = {'.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls', '.csv', '.pdf', '.hwp'}
    
    # 제외할 폴더 패턴
    EXCLUDED_DIRS = {
        '.git', 'node_modules', 'venv', 'env', '__pycache__', 
        '.vscode', '.idea', 'dist', 'build', 'out', 'target',
        '.next', '.nuxt', '.cache', '.temp', '.tmp',
        'vendor', 'packages', 'bower_components'
    }
    
    # 제외할 파일 패턴 (정확한 이름 매칭)
    EXCLUDED_FILES = {
        'desktop.ini', 'thumbs.db', 'Thumbs.db', 'ehthumbs.db',
        '.DS_Store', 'Icon\r', '.gitignore', '.gitattributes'
    }
    
    # Office 임시 파일 패턴 (접두사)
    OFFICE_TEMP_PREFIXES = ('~$', '~WRL')
    
    # 제외할 확장자 (실행파일 및 바이너리)
    EXCLUDED_EXTENSIONS = {
        '.exe', '.dll', '.sys', '.bin', '.so', '.dylib', '.a', '.lib',
        '.o', '.obj', '.class', '.pyc', '.pyo', '.pyd',
        '.iso', '.img', '.dmg', '.vhd', '.vmdk',
        '.zip', '.rar', '.7z', '.tar', '.gz', '.bz2',
        '.mp3', '.mp4', '.avi', '.mkv', '.mov', '.flv',
        '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.ico', '.svg',
        '.ttf', '.otf', '.woff', '.woff2', '.eot'
    }
    
    # 제외할 경로 접두사 (대소문자 구분 없음)
    EXCLUDED_PATH_PREFIXES = [
        'C:\\Windows',
        'C:\\Program Files',
        'C:\\Program Files (x86)',
        'C:\\ProgramData',
        'C:\\$Recycle.Bin',
        'C:\\System Volume Information',
        'C:\\Recovery',
        'C:\\pagefile.sys',
        'C:\\hiberfil.sys',
        'C:\\swapfile.sys'
    ]
    
    def __init__(self, db_manager: DatabaseManager, log_dir: str = None, enable_activity_monitor: bool = True):
        """
        파일 인덱서 초기화
        
        Args:
            db_manager: 데이터베이스 매니저 인스턴스
            log_dir: 로그 파일 디렉토리 (기본: python-backend/logs)
            enable_activity_monitor: 사용자 활동 모니터링 활성화 여부
        """
        self.db = db_manager
        self.is_running = False
        self.current_thread: Optional[threading.Thread] = None
        self.progress_callback: Optional[Callable] = None
        self.log_callback: Optional[Callable] = None
        self.status_callback: Optional[Callable] = None
        self.stop_flag = threading.Event()
        
        # 로그 디렉토리 설정
        if log_dir is None:
            log_dir = os.path.join(os.path.dirname(__file__), 'logs')
        self.log_dir = log_dir
        os.makedirs(self.log_dir, exist_ok=True)
        
        # 로그 파일 경로
        self.skipcheck_file = os.path.join(self.log_dir, 'skipcheck.txt')
        self.error_file = os.path.join(self.log_dir, 'error.txt')
        self.indexing_log_file = os.path.join(self.log_dir, 'indexing_log.txt')  # 통합 인덱싱 로그
        self.indexed_file = os.path.join(self.log_dir, 'Indexed.txt')  # 성공한 인덱싱 결과
        
        # 사용자 정의 제외 패턴
        self.custom_excluded_patterns: List[str] = []
        
        # Skip된 파일 목록 (재시도용)
        self.skipped_files: Dict[str, Dict[str, any]] = {}  # {path: {reason, time, retry_count}}
        self.skipped_files_lock = threading.Lock()
        
        # 재시도 스레드
        self.retry_thread: Optional[threading.Thread] = None
        self.retry_stop_flag = threading.Event()
        self.retry_interval = 300  # 5분 (초 단위)
        
        # 자동 인덱싱 스레드 (주기적 백그라운드 인덱싱)
        self.auto_indexing_thread: Optional[threading.Thread] = None
        self.auto_indexing_stop_flag = threading.Event()
        self.auto_indexing_interval = 1800  # 30분 (초 단위)
        self.auto_indexing_paths: List[str] = []  # 자동 인덱싱할 경로들
        self.is_auto_indexing_enabled = False
        
        # 인덱싱 로그 (메모리, 최근 500개)
        self.indexing_logs: List[Dict[str, str]] = []
        self.indexing_logs_lock = threading.Lock()
        self.max_logs = 500
        
        # 사용자 활동 모니터링 (키보드/마우스 입력 감지)
        self.activity_monitor = UserActivityMonitor(idle_threshold=2.0) if enable_activity_monitor else None
        self.enable_activity_monitor = enable_activity_monitor
        self.paused_count = 0  # 일시정지된 횟수 (통계용)
        
        # 통계
        self.stats = {
            'total_files': 0,
            'indexed_files': 0,
            'skipped_files': 0,
            'error_files': 0,
            'new_files': 0,
            'modified_files': 0,
            'deleted_files': 0,
            'start_time': None,
            'end_time': None,
            'paused_count': 0
        }
    
    def add_exclusion_pattern(self, pattern: str):
        """
        사용자 정의 제외 패턴 추가
        
        Args:
            pattern: 제외할 경로 패턴 (예: "C:\\SecureFolder", "*/private/*")
        """
        if pattern and pattern not in self.custom_excluded_patterns:
            self.custom_excluded_patterns.append(pattern)
            logger.info(f"제외 패턴 추가: {pattern}")
    
    def remove_exclusion_pattern(self, pattern: str):
        """사용자 정의 제외 패턴 제거"""
        if pattern in self.custom_excluded_patterns:
            self.custom_excluded_patterns.remove(pattern)
            logger.info(f"제외 패턴 제거: {pattern}")
    
    def clear_exclusion_patterns(self):
        """모든 사용자 정의 제외 패턴 제거"""
        self.custom_excluded_patterns = []
        logger.info("모든 사용자 정의 제외 패턴 제거")
    
    def get_exclusion_patterns(self) -> List[str]:
        """사용자 정의 제외 패턴 조회"""
        return self.custom_excluded_patterns.copy()
    
    def start_indexing(self, root_paths: List[str], 
                      progress_callback: Optional[Callable] = None,
                      log_callback: Optional[Callable] = None,
                      status_callback: Optional[Callable] = None,
                      silent_mode: bool = False):
        """
        비동기 인덱싱 시작
        
        Args:
            root_paths: 인덱싱할 루트 디렉토리 리스트
            progress_callback: 진행 상황 콜백 (current, total, path)
            log_callback: 로그 콜백 (status, filename, detail)
            status_callback: 상태 콜백 (status_text)
            silent_mode: Silent 모드 (자동 인덱싱용, 로그 최소화)
        """
        if self.is_running:
            logger.warning("인덱싱이 이미 실행 중입니다.")
            return False
        
        self.progress_callback = progress_callback
        self.log_callback = log_callback
        self.status_callback = status_callback
        self.stop_flag.clear()
        
        # Worker 쓰레드에서 실행
        self.current_thread = threading.Thread(
            target=self._indexing_worker,
            args=(root_paths, silent_mode),
            daemon=True
        )
        self.current_thread.start()
        return True
    
    def _write_indexing_log(self, status: str, path: str, detail: str):
        """
        통합 인덱싱 로그 기록 (indexing_log.txt)
        
        Format: [Timestamp] Status | Path | Detail
        
        Args:
            status: 상태 ('Indexing', 'Success', 'Skip', 'Error', 'Retry Success')
            path: 파일 전체 경로
            detail: 상세 정보
        """
        try:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_line = f"[{timestamp}] {status:15s} | {path} | {detail}\n"
            
            with open(self.indexing_log_file, 'a', encoding='utf-8') as f:
                f.write(log_line)
        
        except Exception as e:
            logger.error(f"통합 로그 기록 오류: {e}")
    
    def _write_indexed_file(self, path: str, char_count: int, token_count: int, content: str = None):
        """
        성공한 인덱싱 결과를 Indexed.txt에 기록
        
        Format: 
        ================================================================================
        [Timestamp] 
        디렉토리: {directory}
        파일명: {filename}
        통계: {char_count}자 / {token_count}토큰
        --------------------------------------------------------------------------------
        [인덱스된 내용 미리보기 - 최대 500자]
        {content_preview}
        ================================================================================
        
        Args:
            path: 파일 전체 경로
            char_count: 추출된 문자 수
            token_count: 토큰(단어) 수
            content: 인덱스된 텍스트 내용 (선택사항)
        """
        try:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            directory = os.path.dirname(path)
            filename = os.path.basename(path)
            
            # 로그 엔트리 작성
            log_entry = f"\n{'='*80}\n"
            log_entry += f"[{timestamp}]\n"
            log_entry += f"디렉토리: {directory}\n"
            log_entry += f"파일명: {filename}\n"
            log_entry += f"통계: {char_count:,}자 / {token_count:,}토큰\n"
            log_entry += f"{'-'*80}\n"
            
            # 내용 미리보기 (최대 500자)
            if content:
                preview_length = 500
                content_preview = content[:preview_length]
                if len(content) > preview_length:
                    content_preview += "... (이하 생략)"
                log_entry += f"[인덱스된 내용 미리보기]\n{content_preview}\n"
            else:
                log_entry += f"[내용 없음]\n"
            
            log_entry += f"{'='*80}\n"
            
            with open(self.indexed_file, 'a', encoding='utf-8') as f:
                f.write(log_entry)
        
        except Exception as e:
            logger.error(f"Indexed.txt 기록 오류: {e}")
    
    def _log_skip(self, path: str, reason: str):
        """
        Skip 로그 기록 (skipcheck.txt 및 재시도 목록)
        
        Format: [Timestamp] Path : Reason
        """
        try:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_line = f"[{timestamp}] {path} : {reason}\n"
            
            with open(self.skipcheck_file, 'a', encoding='utf-8') as f:
                f.write(log_line)
            
            # 통합 로그에도 기록
            self._write_indexing_log('Skip', path, reason)
            
            # 재시도 가능한 오류인 경우 목록에 추가
            retryable_reasons = [
                'File locked', 'Permission denied', 'Parsing timeout',
                'Password protected',  # 사용자가 암호 해제할 수 있음
                'File is open'  # 사용자가 파일을 닫으면 재시도
            ]
            
            if any(retryable in reason for retryable in retryable_reasons):
                with self.skipped_files_lock:
                    if path not in self.skipped_files:
                        self.skipped_files[path] = {
                            'reason': reason,
                            'time': time.time(),
                            'retry_count': 0
                        }
            
            # 메모리에 로그 추가
            self._add_log_to_memory('Skip', path, reason)
            
            # UI 로그 콜백
            if self.log_callback:
                self.log_callback('Skip', filename, reason)
        
        except Exception as e:
            logger.error(f"Skip 로그 기록 오류: {e}")
    
    def _log_error(self, path: str, error: Exception):
        """
        에러 로그 기록 (error.txt)
        
        트레이스백 포함
        """
        try:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            error_msg = f"\n[{timestamp}] {path}\n"
            error_msg += f"Error: {str(error)}\n"
            error_msg += f"Traceback:\n{traceback.format_exc()}\n"
            error_msg += "=" * 80 + "\n"
            
            with open(self.error_file, 'a', encoding='utf-8') as f:
                f.write(error_msg)
            
            # 통합 로그에도 기록
            self._write_indexing_log('Error', path, f"Error: {str(error)}")
            
            # 메모리에 로그 추가
            self._add_log_to_memory('Error', path, str(error))
            
            # UI 로그 콜백
            if self.log_callback:
                self.log_callback('Error', filename, str(error))
        
        except Exception as e:
            logger.error(f"에러 로그 기록 오류: {e}")
    
    def _add_log_to_memory(self, status: str, path: str, detail: str):
        """
        메모리에 로그 추가 (API 조회용)
        
        Args:
            status: 상태 ('Success', 'Skip', 'Error', 'Indexing', 'Retry Success')
            path: 파일 전체 경로
            detail: 상세 정보
        """
        with self.indexing_logs_lock:
            log_entry = {
                'time': datetime.now().strftime('%H:%M:%S'),
                'status': status,
                'path': path,  # 전체 경로 저장
                'filename': os.path.basename(path),  # 파일명도 별도로 저장
                'detail': detail
            }
            self.indexing_logs.insert(0, log_entry)  # 최신 로그를 앞에
            # 최대 개수 유지
            if len(self.indexing_logs) > self.max_logs:
                self.indexing_logs = self.indexing_logs[:self.max_logs]
    
    def get_recent_logs(self, count: int = 100) -> List[Dict[str, str]]:
        """
        최근 로그 조회
        
        Args:
            count: 조회할 로그 수
        
        Returns:
            로그 리스트
        """
        with self.indexing_logs_lock:
            return self.indexing_logs[:count]
    
    def clear_logs(self):
        """로그 초기화"""
        with self.indexing_logs_lock:
            self.indexing_logs = []
    
    def _add_to_retry_queue(self, file_path: str, reason: str):
        """
        재시도 큐에 파일 추가
        
        Args:
            file_path: 파일 경로
            reason: Skip 사유
        """
        with self.skipped_files_lock:
            if file_path not in self.skipped_files:
                self.skipped_files[file_path] = {
                    'reason': reason,
                    'time': time.time(),
                    'retry_count': 0
                }
                logger.info(f"재시도 큐 추가: {file_path} (사유: {reason})")
    
    def _count_tokens(self, text: str) -> int:
        """
        텍스트의 토큰(단어) 수 계산 (다국어 지원)
        
        - 한글/중국어/일본어(CJK): 각 문자를 1토큰으로 계산
        - 영어/기타 언어: 공백 기준 단어로 계산
        
        Args:
            text: 텍스트 문자열
        
        Returns:
            토큰 수
        """
        if not text:
            return 0
        
        token_count = 0
        
        # CJK 문자 패턴 (한글, 중국어, 일본어)
        # - 한글: \uAC00-\uD7AF (가-힣), \u1100-\u11FF, \u3130-\u318F
        # - 중국어: \u4E00-\u9FFF
        # - 일본어: \u3040-\u309F (히라가나), \u30A0-\u30FF (가타카나)
        cjk_pattern = re.compile(r'[\uAC00-\uD7AF\u1100-\u11FF\u3130-\u318F\u4E00-\u9FFF\u3040-\u309F\u30A0-\u30FF]')
        
        # 텍스트를 CJK 문자와 비CJK 부분으로 분리
        parts = re.split(r'([\uAC00-\uD7AF\u1100-\u11FF\u3130-\u318F\u4E00-\u9FFF\u3040-\u309F\u30A0-\u30FF]+)', text)
        
        for part in parts:
            if not part.strip():
                continue
            
            # CJK 문자인 경우 각 문자를 1토큰으로
            if cjk_pattern.search(part):
                # 공백 제거 후 문자 수 계산
                token_count += len(part.strip())
            else:
                # 영어 등 공백 기반 언어는 단어 수로 계산
                words = part.split()
                token_count += len(words)
        
        return token_count
    
    def _log_success(self, path: str, char_count: int, token_count: int = 0, db_saved: bool = True, content: str = None):
        """
        성공 로그 (UI + 파일)
        
        Args:
            path: 파일 경로
            char_count: 추출된 문자 수
            token_count: 토큰(단어) 수
            db_saved: DB 저장 성공 여부
            content: 인덱스된 텍스트 내용
        """
        filename = os.path.basename(path)
        db_status = "✓ DB 저장 완료" if db_saved else "⊗ DB 저장 대기"
        detail = f'{char_count:,}자 / {token_count:,}토큰 | {db_status}'
        
        # 통합 로그에 기록
        self._write_indexing_log('Success', path, detail)
        
        # DB 저장이 완료된 경우에만 Indexed.txt에 기록
        if db_saved:
            self._write_indexed_file(path, char_count, token_count, content)
        
        # 메모리에 로그 추가
        self._add_log_to_memory('Success', path, detail)
        
        # UI 콜백
        if self.log_callback:
            self.log_callback('Success', filename, detail)
    
    def _log_indexing(self, path: str):
        """
        인덱싱 시작 로그 (현재 처리 중인 파일 표시)
        
        Args:
            path: 파일 경로
        """
        filename = os.path.basename(path)
        directory = os.path.dirname(path)
        # 디렉토리 경로를 상대적으로 표시 (너무 길면 생략)
        if len(directory) > 60:
            directory = '...' + directory[-57:]
        detail = f'처리 중... [{directory}]'
        
        # 통합 로그에 기록
        self._write_indexing_log('Indexing', path, detail)
        
        # 메모리에 로그 추가
        self._add_log_to_memory('Indexing', path, detail)
        
        # UI 콜백
        if self.log_callback:
            self.log_callback('Indexing', filename, detail)
    
    def _update_status(self, status: str):
        """상태 업데이트"""
        if self.status_callback:
            self.status_callback(status)
    
    def stop_indexing(self):
        """인덱싱 중지 - 2중 체크로 확실하게 중단"""
        if not self.is_running:
            logger.info("인덱싱이 실행 중이 아닙니다.")
            return
        
        logger.info("========================================")
        logger.info("인덱싱 중지 요청 - 1차 시도")
        logger.info("========================================")
        
        # 1차: stop_flag 설정
        self.stop_flag.set()
        
        # 재시도 워커도 함께 중지
        if self.retry_thread and self.retry_thread.is_alive():
            logger.info("재시도 워커 중지 시도...")
            self.retry_stop_flag.set()
        
        # 워커 스레드가 종료될 때까지 대기 (최대 10초)
        if self.current_thread and self.current_thread.is_alive():
            logger.info("워커 스레드 종료 대기 중...")
            self.current_thread.join(timeout=10)
            
            # 1차 종료 확인
            if self.current_thread.is_alive():
                logger.warning("⚠ 1차 중지 실패 - 워커 스레드가 아직 실행 중")
                logger.info("========================================")
                logger.info("인덱싱 중지 요청 - 2차 시도 (강제)")
                logger.info("========================================")
                
                # 2차: 강제 종료 시도 - stop_flag 재설정
                self.stop_flag.set()
                self.is_running = False  # 강제로 상태 변경
                
                # 추가 대기 (5초)
                self.current_thread.join(timeout=5)
                
                # 2차 종료 확인
                if self.current_thread.is_alive():
                    logger.error("❌ 2차 중지 실패 - 워커 스레드 강제 종료 불가")
                    logger.error("❌ 시스템 재시작이 필요할 수 있습니다")
                else:
                    logger.info("✓ 2차 시도로 인덱싱 중지 완료")
            else:
                logger.info("✓ 1차 시도로 인덱싱 중지 완료")
        
        # 재시도 워커 종료 확인
        if self.retry_thread and self.retry_thread.is_alive():
            self.retry_thread.join(timeout=2)
            if self.retry_thread.is_alive():
                logger.warning("⚠ 재시도 워커 중지 실패")
            else:
                logger.info("✓ 재시도 워커 중지 완료")
        
        # 최종 상태 확인
        logger.info("========================================")
        logger.info(f"최종 상태: is_running={self.is_running}, 워커 살아있음={self.current_thread.is_alive() if self.current_thread else False}")
        logger.info("========================================")
    
    def stop_retry_worker(self):
        """재시도 워커 중지"""
        if self.retry_thread and self.retry_thread.is_alive():
            logger.info("재시도 워커 중지 요청...")
            self.retry_stop_flag.set()
            self.retry_thread.join(timeout=2)
            logger.info("재시도 워커 중지됨")
    
    # ========================================
    # 자동 인덱싱 (주기적 백그라운드)
    # ========================================
    
    def start_auto_indexing(self, paths: List[str], interval_minutes: int = 30):
        """
        자동 인덱싱 시작 (주기적으로 변경사항 자동 반영)
        
        Args:
            paths: 자동 인덱싱할 경로 목록
            interval_minutes: 인덱싱 주기 (분 단위, 기본: 30분)
        """
        if self.is_auto_indexing_enabled:
            logger.warning("자동 인덱싱이 이미 실행 중입니다.")
            return
        
        self.auto_indexing_paths = paths
        self.auto_indexing_interval = interval_minutes * 60  # 분 -> 초
        self.auto_indexing_stop_flag.clear()
        self.is_auto_indexing_enabled = True
        
        self.auto_indexing_thread = threading.Thread(target=self._auto_indexing_worker, daemon=True)
        self.auto_indexing_thread.start()
        
        logger.info(f"🤖 자동 인덱싱 시작 (주기: {interval_minutes}분)")
        logger.info(f"   감시 경로: {paths}")
    
    def stop_auto_indexing(self):
        """자동 인덱싱 중지"""
        if not self.is_auto_indexing_enabled:
            return
        
        logger.info("자동 인덱싱 중지 요청...")
        self.auto_indexing_stop_flag.set()
        self.is_auto_indexing_enabled = False
        
        if self.auto_indexing_thread and self.auto_indexing_thread.is_alive():
            self.auto_indexing_thread.join(timeout=3)
        
        logger.info("자동 인덱싱 중지됨")
    
    def _auto_indexing_worker(self):
        """자동 인덱싱 워커 스레드 - 주기적으로 변경사항 감지 및 인덱싱"""
        logger.info("🤖 자동 인덱싱 워커 시작")
        
        # 첫 실행은 5분 후부터 시작
        initial_delay = 300  # 5분
        
        if self._wait_with_interrupt(initial_delay):
            return  # 중지 요청됨
        
        while not self.auto_indexing_stop_flag.is_set():
            try:
                # 수동 인덱싱이 실행 중이면 대기
                if self.is_running:
                    logger.debug("🤖 자동 인덱싱: 수동 인덱싱 실행 중 - 대기")
                    if self._wait_with_interrupt(60):  # 1분 대기
                        break
                    continue
                
                logger.info("========================================")
                logger.info("🤖 자동 인덱싱 시작 (Silent Mode)")
                logger.info("========================================")
                
                # 삭제된 파일 정리
                self._cleanup_deleted_files()
                
                # 증분 인덱싱 실행 (변경/추가 파일만)
                if self.auto_indexing_paths:
                    self.start_indexing(self.auto_indexing_paths, silent_mode=True)
                    
                    # 인덱싱 완료 대기
                    while self.is_running and not self.auto_indexing_stop_flag.is_set():
                        time.sleep(1)
                
                logger.info("🤖 자동 인덱싱 완료 - 다음 실행까지 대기")
                
                # 다음 주기까지 대기
                if self._wait_with_interrupt(self.auto_indexing_interval):
                    break
                
            except Exception as e:
                logger.error(f"자동 인덱싱 워커 오류: {e}")
                logger.error(traceback.format_exc())
                # 오류 발생 시 5분 대기 후 재시도
                if self._wait_with_interrupt(300):
                    break
        
        logger.info("🤖 자동 인덱싱 워커 종료")
    
    def _wait_with_interrupt(self, seconds: int) -> bool:
        """
        지정된 시간 동안 대기, 중간에 중지 플래그 확인
        
        Returns:
            True면 중지 요청됨, False면 정상 대기 완료
        """
        elapsed = 0
        check_interval = 1  # 1초마다 체크
        
        while elapsed < seconds:
            if self.auto_indexing_stop_flag.is_set():
                return True
            time.sleep(min(check_interval, seconds - elapsed))
            elapsed += check_interval
        
        return False
    
    def _cleanup_deleted_files(self):
        """DB에서 삭제된 파일 제거"""
        try:
            logger.info("🗑️ 삭제된 파일 정리 시작...")
            
            # DB에서 모든 인덱싱된 파일 경로 가져오기
            indexed_files = self.db.get_all_indexed_file_paths()
            
            deleted_count = 0
            for file_path in indexed_files:
                # 파일이 실제로 존재하는지 확인
                if not os.path.exists(file_path):
                    logger.info(f"🗑️ 삭제된 파일 제거: {file_path}")
                    self.db.delete_file(file_path)
                    deleted_count += 1
            
            if deleted_count > 0:
                logger.info(f"✅ 삭제된 파일 {deleted_count}개 제거 완료")
                self.stats['deleted_files'] += deleted_count
            else:
                logger.debug("✅ 삭제된 파일 없음")
        
        except Exception as e:
            logger.error(f"삭제된 파일 정리 오류: {e}")
            logger.error(traceback.format_exc())
    
    def cleanup(self):
        """인덱서 리소스 정리 및 Lock 해제 - 강화된 종료 보장"""
        logger.info("========================================")
        logger.info("인덱서 리소스 정리 시작 (강화 모드)")
        logger.info("========================================")
        
        try:
            # 0단계: 자동 인덱싱 및 사용자 활동 모니터 중지
            if self.is_auto_indexing_enabled:
                logger.info("0단계-1: 자동 인덱싱 중지...")
                self.stop_auto_indexing()
            
            if self.activity_monitor:
                logger.info("0단계-2: 사용자 활동 모니터 중지...")
                self.activity_monitor.stop()
            
            # 1단계: stop_flag 설정
            self.stop_flag.set()
            self.retry_stop_flag.set()
            self.auto_indexing_stop_flag.set()
            logger.info("1단계: 모든 stop_flag 설정 완료")
            
            # 2단계: 재시도 워커 강제 종료 (우선 처리)
            if self.retry_thread and self.retry_thread.is_alive():
                logger.info("2단계: 재시도 워커 종료 시도...")
                self.retry_stop_flag.set()
                self.retry_thread.join(timeout=3)
                
                if self.retry_thread.is_alive():
                    logger.warning("⚠ 재시도 워커 3초 내 종료 실패 - 강제 종료")
                    # 추가 대기
                    self.retry_thread.join(timeout=2)
                else:
                    logger.info("✓ 재시도 워커 종료 완료")
            
            # 3단계: 메인 인덱싱 스레드 종료 (다중 시도)
            if self.is_running or (self.current_thread and self.current_thread.is_alive()):
                logger.info("3단계: 메인 인덱싱 스레드 종료 시도...")
                
                # 1차 시도: 정상 종료 (10초 대기)
                self.is_running = False  # 상태 플래그 강제 변경
                if self.current_thread and self.current_thread.is_alive():
                    logger.info("  - 1차 시도: 정상 종료 대기 (10초)...")
                    self.current_thread.join(timeout=10)
                
                # 2차 시도: 추가 대기 (5초)
                if self.current_thread and self.current_thread.is_alive():
                    logger.warning("  - 1차 시도 실패, 2차 시도 (5초)...")
                    self.current_thread.join(timeout=5)
                
                # 최종 확인
                if self.current_thread and self.current_thread.is_alive():
                    logger.error("  ⚠ 메인 인덱싱 스레드 종료 실패 (15초 초과)")
                else:
                    logger.info("  ✓ 메인 인덱싱 스레드 종료 완료")
            
            # 4단계: 메모리 정리
            logger.info("4단계: 메모리 정리...")
            with self.skipped_files_lock:
                self.skipped_files.clear()
            
            with self.indexing_logs_lock:
                self.indexing_logs.clear()
            
            # 5단계: 최종 상태 확인
            main_alive = self.current_thread.is_alive() if self.current_thread else False
            retry_alive = self.retry_thread.is_alive() if self.retry_thread else False
            
            logger.info("========================================")
            logger.info("인덱서 리소스 정리 완료")
            logger.info(f"  - is_running: {self.is_running}")
            logger.info(f"  - 메인 스레드 살아있음: {main_alive}")
            logger.info(f"  - 재시도 워커 살아있음: {retry_alive}")
            logger.info("========================================")
            
        except Exception as e:
            logger.error(f"❌ 인덱서 정리 중 오류: {e}")
            logger.error(traceback.format_exc())
    
    def _indexing_worker(self, root_paths: List[str], silent_mode: bool = False):
        """인덱싱 워커 (백그라운드 쓰레드) - 증분 색인"""
        self.is_running = True
        self.silent_mode = silent_mode  # Silent 모드 플래그 저장
        self.stats = {
            'total_files': 0,
            'indexed_files': 0,
            'skipped_files': 0,
            'error_files': 0,
            'new_files': 0,
            'modified_files': 0,
            'deleted_files': 0,
            'start_time': time.time(),
            'end_time': None,
            'paused_count': 0
        }
        
        # 사용자 활동 모니터링 시작
        if self.activity_monitor and self.enable_activity_monitor:
            self.activity_monitor.start()
            logger.info("사용자 활동 모니터링 활성화: 키보드/마우스 입력 감지 시 자동 일시정지")
        
        logger.info(f"인덱싱 시작: {root_paths}")
        self._update_status("파일 수집 중...")
        
        try:
            # 1단계: 파일 목록 수집
            all_files = []
            for root_path in root_paths:
                if self.stop_flag.is_set():
                    break
                all_files.extend(self._collect_files(root_path))
            
            self.stats['total_files'] = len(all_files)
            logger.info(f"수집된 파일: {len(all_files)}개")
            self._update_status(f"총 {len(all_files)}개 파일 발견")
            
            # 2단계: 증분 인덱싱 (New/Modified 파일만)
            self._update_status("증분 인덱싱 중...")
            self._process_files_incremental(all_files)
            
            # 3단계: 삭제된 파일 정리
            self._update_status("삭제된 파일 정리 중...")
            self._cleanup_deleted_files(all_files)
            
            # 4단계: DB 최적화 (VACUUM)
            self._update_status("데이터베이스 최적화 중...")
            self.db.vacuum()
            
            # 인덱스 최적화
            self.db.optimize()
            
        except Exception as e:
            error_msg = f"인덱싱 워커 오류: {type(e).__name__} - {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            self._log_error("IndexingWorker", e)
            
            # UI에 에러 원인 표시
            if self.log_callback:
                self.log_callback('Error', '인덱싱 중단', f'{type(e).__name__}: {str(e)}')
        
        finally:
            # 사용자 활동 모니터링 중지
            if self.activity_monitor:
                self.activity_monitor.stop()
            
            self.stats['end_time'] = time.time()
            elapsed = self.stats['end_time'] - self.stats['start_time']
            summary = f"완료: {self.stats['indexed_files']}개 인덱싱 ({elapsed:.2f}초)"
            if self.stats['paused_count'] > 0:
                summary += f" | 일시정지 {self.stats['paused_count']}회"
            logger.info(summary)
            self._update_status(summary)
            self.is_running = False
            
            # 재시도 워커 시작 (Skip된 파일이 있는 경우)
            with self.skipped_files_lock:
                if self.skipped_files:
                    logger.info(f"재시도 워커 시작: Skip된 파일 {len(self.skipped_files)}개")
                    self.start_retry_worker()
    
    def _process_files_incremental(self, all_files: List[str]):
        """증분 파일 처리 (New/Modified만) - 리소스 사용 최소화"""
        batch_size = 2  # 2개 파일마다 DB Commit (리소스 절약)
        batch = []
        last_progress_time = time.time()
        stall_warning_threshold = 120  # 2분 동안 진행 없으면 경고
        file_delay = 0.01  # 파일 처리 간 0.01초 지연 (즉각적인 활동 감지)
        
        for i, file_path in enumerate(all_files):
            if self.stop_flag.is_set():
                logger.info("인덱싱 중지됨 (사용자 요청)")
                if self.log_callback:
                    self.log_callback('Info', '인덱싱 중지', '사용자가 중지를 요청했습니다')
                break
            
            # 사용자 활동 체크 - 키보드/마우스 입력 감지 시 대기
            if self.activity_monitor and self.enable_activity_monitor:
                if self.activity_monitor.is_user_active():
                    # 사용자 활동 감지 - 즉시 일시정지
                    self.stats['paused_count'] += 1
                    idle_time = self.activity_monitor.get_idle_time()
                    logger.info(f"⏸️ 사용자 활동 감지 (즉시 중단) - 2초 대기 중...")
                    self._update_status(f"⏸️ 사용자 작업 중 - 2초 대기 중...")
                    
                    # UI 로그
                    if self.log_callback:
                        self.log_callback('Info', '일시정지', '⏸️ 사용자 작업 중 - 2초 대기')
                    
                    # 유휴 상태가 될 때까지 대기 (2초 동안 입력 없을 때까지)
                    # check_interval을 0.1초로 줄여서 더 즉각적으로 반응
                    if not self.activity_monitor.wait_until_idle(check_interval=0.1, stop_flag=self.stop_flag):
                        # 중지 요청됨
                        break
                    
                    # 재개
                    logger.info("▶️ 사용자 활동 없음 (2초 경과) - 인덱싱 재개")
                    self._update_status("▶️ 인덱싱 재개 중...")
                    if self.log_callback:
                        self.log_callback('Info', '재개', '▶️ 인덱싱 재개됨')
            
            # 진행 상황 체크 (2분 이상 멈춤 감지)
            current_time = time.time()
            if current_time - last_progress_time > stall_warning_threshold:
                warning_msg = f"⚠ 인덱싱 진행 지연 감지: {file_path} 처리 중 {stall_warning_threshold}초 경과"
                logger.warning(warning_msg)
                if self.log_callback:
                    self.log_callback('Error', '진행 지연', f'{os.path.basename(file_path)} 처리 중 지연')
                last_progress_time = current_time
            
            try:
                # 진행 상황 콜백
                if self.progress_callback:
                    self.progress_callback(i + 1, len(all_files), file_path)
                
                # 파일 잠금 체크 제거 - 임시 파일 복사로 처리하므로 불필요
                # 각 파일 타입의 extract 함수가 _copy_to_temp를 사용하여
                # 사용자가 열어둔 파일도 안전하게 인덱싱합니다
                
                # 파일 크기 체크 (100MB 초과 시 스킵)
                try:
                    file_size = os.path.getsize(file_path)
                    if file_size > MAX_FILE_SIZE:
                        self._log_skip(file_path, f"Size exceeded ({file_size / 1024 / 1024:.1f}MB)")
                        self.stats['skipped_files'] += 1
                        continue
                except Exception:
                    pass
                
                # 증분 인덱싱: New or Modified?
                current_mtime = os.path.getmtime(file_path)
                indexed_mtime = self.db.get_file_mtime(file_path)
                
                if indexed_mtime is not None:
                    # 파일이 이미 인덱싱됨
                    if abs(current_mtime - indexed_mtime) < 1.0:
                        # 수정되지 않음 - 스킵
                        self.stats['skipped_files'] += 1
                        continue
                    else:
                        # 수정됨 - 재인덱싱
                        is_new = False
                        self.stats['modified_files'] += 1
                else:
                    # 새 파일
                    is_new = True
                    self.stats['new_files'] += 1
                
                # 중지 요청 체크 (파일 처리 전)
                if self.stop_flag.is_set():
                    logger.info("인덱싱 중지됨 (사용자 요청 - 파일 처리 전)")
                    if self.log_callback:
                        self.log_callback('Info', '인덱싱 중지', '사용자가 중지를 요청했습니다')
                    break
                
                # 현재 처리 중인 파일 로그
                self._log_indexing(file_path)
                
                # 텍스트 추출 (타임아웃 체크)
                content = self._extract_text_safe(file_path)
                
                if content:
                    # 토큰 수 계산
                    token_count = self._count_tokens(content)
                    
                    if is_new:
                        # 새 파일은 배치에 추가 (로그는 DB 저장 완료 후 생성)
                        batch.append((file_path, content, current_mtime, token_count))
                        self.stats['indexed_files'] += 1
                    else:
                        # 수정된 파일은 즉시 업데이트
                        try:
                            self.db.update_file(file_path, content, current_mtime)
                            self.stats['indexed_files'] += 1
                            # DB 저장 완료 상태로 로그
                            self._log_success(file_path, len(content), token_count, db_saved=True, content=content)
                        except Exception as e:
                            logger.error(f"DB 업데이트 오류 [{file_path}]: {e}")
                            self._log_error(file_path, e)
                            self.stats['error_files'] += 1
                            continue
                    
                    # 배치가 가득 찼으면 DB에 저장
                    if len(batch) >= batch_size:
                        # 중지 요청 체크 (배치 저장 전)
                        if self.stop_flag.is_set():
                            logger.info("인덱싱 중지됨 (사용자 요청 - 배치 저장 전)")
                            if self.log_callback:
                                self.log_callback('Info', '인덱싱 중지', '사용자가 중지를 요청했습니다')
                            break
                        
                        try:
                            # 배치 저장 (토큰 수 제외)
                            batch_for_db = [(path, content, mtime) for path, content, mtime, _ in batch]
                            self.db.insert_files_batch(batch_for_db)
                            # 배치 저장 후 DB 저장 완료 로그 생성
                            for saved_path, saved_content, _, saved_token_count in batch:
                                # 3. DB 저장 완료 로그
                                self._log_success(saved_path, len(saved_content), saved_token_count, db_saved=True, content=saved_content)
                            batch = []
                            last_progress_time = time.time()  # 진행 시간 업데이트
                            # 배치 저장 후 지연 (IO 부하 감소)
                            time.sleep(0.5)
                        except Exception as e:
                            logger.error(f"DB 배치 저장 오류: {e}")
                            if self.log_callback:
                                self.log_callback('Error', 'DB 저장', f'배치 저장 오류: {str(e)}')
                            batch = []
                else:
                    self.stats['skipped_files'] += 1
                
                # 파일 처리 간 지연 (CPU/IO 부하 감소)
                time.sleep(file_delay)
            
            except PermissionError as e:
                # 파일 잠금
                self._log_skip(file_path, "File locked or Permission denied")
                self.stats['skipped_files'] += 1
            
            except Exception as e:
                error_type = type(e).__name__
                error_msg = str(e)
                logger.error(f"파일 처리 오류 [{file_path}]: {error_type} - {error_msg}")
                logger.error(f"상세 정보: {traceback.format_exc()}")
                
                # UI에 에러 원인 표시
                self._log_error(file_path, f"{error_type}: {error_msg}")
                self.stats['error_files'] += 1
                
                # 타임아웃 에러인 경우 특별히 표시
                if 'timeout' in error_msg.lower() or error_type == 'TimeoutError':
                    if self.log_callback:
                        self.log_callback('Error', os.path.basename(file_path), f'⏱ 타임아웃 (60초 초과)')
                elif 'memory' in error_msg.lower():
                    if self.log_callback:
                        self.log_callback('Error', os.path.basename(file_path), f'💾 메모리 부족')
                elif error_type == 'PermissionError':
                    if self.log_callback:
                        self.log_callback('Error', os.path.basename(file_path), f'🔒 권한 오류')
        
        # 남은 배치 저장
        if batch:
            try:
                logger.info(f"최종 배치 저장 중: {len(batch)}개 파일")
                # 배치 저장 (토큰 수 제외)
                batch_for_db = [(path, content, mtime) for path, content, mtime, _ in batch]
                self.db.insert_files_batch(batch_for_db)
                # 배치 저장 후 DB 저장 완료 로그 생성
                for saved_path, saved_content, _, saved_token_count in batch:
                    self._log_success(saved_path, len(saved_content), saved_token_count, db_saved=True, content=saved_content)
                logger.info(f"✓ 최종 배치 저장 완료: {len(batch)}개 파일")
            except Exception as e:
                error_msg = f"DB 최종 배치 저장 오류: {type(e).__name__} - {str(e)}"
                logger.error(error_msg)
                logger.error(traceback.format_exc())
                if self.log_callback:
                    self.log_callback('Error', '최종 DB 저장', f'배치 저장 실패: {str(e)}')
    
    def _cleanup_deleted_files(self, current_files: List[str]):
        """삭제된 파일을 DB에서 제거"""
        try:
            # 현재 파일 세트
            current_file_set = set(current_files)
            
            # DB의 모든 파일 경로 조회
            indexed_files = self.db.get_all_indexed_paths()
            
            # 삭제된 파일 찾기
            deleted_files = [f for f in indexed_files if f not in current_file_set]
            
            # 삭제
            for file_path in deleted_files:
                self.db.delete_file(file_path)
                self.stats['deleted_files'] += 1
            
            if deleted_files:
                logger.info(f"삭제된 파일 {len(deleted_files)}개 정리 완료")
        
        except Exception as e:
            logger.error(f"삭제된 파일 정리 오류: {e}")
    
    def _extract_text_safe(self, file_path: str) -> Optional[str]:
        """
        안전한 텍스트 추출 (타임아웃, 예외 처리)
        
        Returns:
            추출된 텍스트 또는 None
        """
        try:
            # 타임아웃 적용 (60초)
            @with_timeout(PARSE_TIMEOUT)
            def extract():
                return self._extract_text(file_path)
            
            return extract()
        
        except TimeoutError:
            self._log_skip(file_path, f"Parsing timeout (>{PARSE_TIMEOUT}s)")
            # 재시도 목록에 추가
            self._add_to_retry_queue(file_path, f"Parsing timeout (>{PARSE_TIMEOUT}s)")
            return None
        
        except PermissionError as e:
            # 파일이 다른 프로그램에서 열려있음
            self._log_skip(file_path, "File is open in another program")
            # 재시도 목록에 추가 (나중에 파일이 닫히면 다시 시도)
            self._add_to_retry_queue(file_path, "File is open in another program")
            return None
        
        except Exception as e:
            # 암호화된 파일, 손상된 파일 등
            error_msg = str(e).lower()
            if 'password' in error_msg or 'encrypted' in error_msg:
                self._log_skip(file_path, "Password protected")
                # 암호 보호 파일도 재시도 목록에 추가 (사용자가 암호 해제할 수 있음)
                self._add_to_retry_queue(file_path, "Password protected")
            elif 'corrupt' in error_msg or 'damaged' in error_msg:
                self._log_skip(file_path, "File corrupted")
                # 손상된 파일은 재시도하지 않음
            else:
                self._log_skip(file_path, f"Parse error: {str(e)[:100]}")
            return None
    
    def _collect_files(self, root_path: str) -> List[str]:
        """
        파일 시스템 크롤링 (제외 규칙 적용)
        
        Args:
            root_path: 루트 디렉토리
        
        Returns:
            파일 경로 리스트
        """
        files = []
        
        try:
            for dirpath, dirnames, filenames in os.walk(root_path):
                # 제외할 디렉토리 필터링
                dirnames[:] = [d for d in dirnames if self._should_include_dir(d, dirpath)]
                
                if self.stop_flag.is_set():
                    break
                
                for filename in filenames:
                    file_path = os.path.join(dirpath, filename)
                    
                    # 파일 포함 여부 확인
                    if not self._should_include_file(filename, file_path):
                        continue
                    
                    files.append(file_path)
        
        except Exception as e:
            logger.error(f"파일 수집 오류 [{root_path}]: {e}")
        
        return files
    
    def _should_include_dir(self, dirname: str, dirpath: str) -> bool:
        """
        디렉토리 포함 여부 확인
        
        Args:
            dirname: 디렉토리 이름
            dirpath: 디렉토리 전체 경로
        
        Returns:
            True면 포함, False면 제외
        """
        # 특수 문자로 시작하는 디렉토리 제외
        if not self._is_valid_name(dirname):
            return False
        
        # 제외 디렉토리 목록에 있으면 제외
        if dirname.lower() in self.EXCLUDED_DIRS:
            return False
        
        # 전체 경로가 제외 경로 접두사에 해당하면 제외
        full_path = os.path.join(dirpath, dirname)
        for excluded_prefix in self.EXCLUDED_PATH_PREFIXES:
            if full_path.lower().startswith(excluded_prefix.lower()):
                return False
        
        return True
    
    def _should_include_file(self, filename: str, filepath: str) -> bool:
        """
        파일 포함 여부 확인
        
        Args:
            filename: 파일 이름
            filepath: 파일 전체 경로
        
        Returns:
            True면 포함, False면 제외
        """
        # 특수 문자로 시작하는 파일 제외
        if not self._is_valid_name(filename):
            return False
        
        # Office 임시 파일 제외 (~$, ~WRL)
        for prefix in self.OFFICE_TEMP_PREFIXES:
            if filename.startswith(prefix):
                return False
        
        # 제외 파일 목록에 있으면 제외
        if filename in self.EXCLUDED_FILES:
            return False
        
        # 확장자 확인
        ext = Path(filepath).suffix.lower()
        
        # 제외 확장자면 제외
        if ext in self.EXCLUDED_EXTENSIONS:
            return False
        
        # 지원하는 확장자가 아니면 제외
        if ext not in self.SUPPORTED_TEXT_EXTENSIONS and ext not in self.SUPPORTED_DOC_EXTENSIONS:
            return False
        
        # 전체 경로가 제외 경로 접두사에 해당하면 제외
        for excluded_prefix in self.EXCLUDED_PATH_PREFIXES:
            if filepath.lower().startswith(excluded_prefix.lower()):
                return False
        
        # 사용자 정의 제외 패턴 체크
        for pattern in self.custom_excluded_patterns:
            # 간단한 패턴 매칭 (와일드카드 지원)
            if self._match_pattern(filepath, pattern):
                return False
        
        return True
    
    def _match_pattern(self, filepath: str, pattern: str) -> bool:
        """
        경로가 패턴과 매칭되는지 확인
        
        Args:
            filepath: 파일 경로
            pattern: 패턴 (와일드카드 * 지원)
        
        Returns:
            True면 매칭됨
        """
        import fnmatch
        
        # 대소문자 구분 없이 매칭
        return fnmatch.fnmatch(filepath.lower(), pattern.lower())
    
    def _is_valid_name(self, name: str) -> bool:
        """특수 문자로 시작하는 파일/폴더 필터링"""
        if not name:
            return False
        return name[0].isalnum() or ord(name[0]) >= 0xAC00  # 영문, 숫자, 한글
    
    def _copy_to_temp(self, file_path: str) -> Optional[str]:
        """
        파일을 임시 폴더에 복사 (강제 복사)
        
        🔥 핵심: 사용자가 열어서 사용 중인 파일도 강제로 읽기 전용 복사!
        임시 복사본을 만들어서 인덱싱하므로 원본 파일은 절대 건드리지 않습니다.
        
        Args:
            file_path: 원본 파일 경로
        
        Returns:
            임시 파일 경로 또는 None (복사 실패 시)
        """
        try:
            # 임시 디렉토리 생성
            temp_dir = tempfile.mkdtemp(prefix='indexer_')
            
            # 파일명 가져오기
            filename = os.path.basename(file_path)
            temp_file_path = os.path.join(temp_dir, filename)
            
            # 방법 1: shutil.copy2 시도 (가장 빠름)
            try:
                shutil.copy2(file_path, temp_file_path)
                logger.debug(f"✅ 임시 파일 복사 완료 (shutil): {filename}")
                return temp_file_path
            except (PermissionError, IOError) as e:
                # 방법 2: 바이너리 읽기 모드로 직접 복사 (더 강력)
                logger.debug(f"shutil 복사 실패, 직접 읽기 시도: {filename}")
                try:
                    with open(file_path, 'rb') as src:
                        data = src.read()
                    with open(temp_file_path, 'wb') as dst:
                        dst.write(data)
                    logger.debug(f"✅ 임시 파일 복사 완료 (직접 읽기): {filename}")
                    return temp_file_path
                except Exception as e2:
                    logger.info(f"⛔ 파일 복사 완전 실패 - Skip: {filename} (원인: {e2})")
                    # 임시 디렉토리 정리
                    try:
                        shutil.rmtree(temp_dir, ignore_errors=True)
                    except:
                        pass
                    return None
            
        except Exception as e:
            logger.debug(f"임시 파일 복사 실패 [{file_path}]: {e}")
            return None
    
    def _cleanup_temp(self, temp_file_path: str):
        """
        임시 파일 및 폴더 정리
        
        Args:
            temp_file_path: 임시 파일 경로
        """
        try:
            if temp_file_path and os.path.exists(temp_file_path):
                # 임시 디렉토리 전체 삭제
                temp_dir = os.path.dirname(temp_file_path)
                shutil.rmtree(temp_dir, ignore_errors=True)
                logger.debug(f"임시 파일 정리 완료: {os.path.basename(temp_file_path)}")
        except Exception as e:
            logger.debug(f"임시 파일 정리 오류: {e}")
    
    def _is_file_locked(self, file_path: str) -> bool:
        """
        파일이 다른 프로세스에서 사용 중인지 확인 (최강화 버전)
        
        사용자가 열어서 사용 중인 파일은 절대 건드리지 않습니다!
        
        Args:
            file_path: 파일 경로
        
        Returns:
            True면 파일이 잠겨있음 (사용자가 열어서 사용 중) - 절대 건드리지 않음!
        """
        import errno
        
        # 방법 1: 독점 쓰기 모드로 열기 시도
        try:
            # 'r+b' 모드로 열기 시도 - 다른 프로그램이 독점 모드로 열었으면 실패
            file_handle = open(file_path, 'r+b')
            file_handle.close()
        except PermissionError:
            # 권한 없거나 다른 프로그램이 독점 사용 중
            logger.debug(f"파일 잠금 감지 (PermissionError): {file_path}")
            return True
        except IOError as e:
            # 파일이 사용 중인 경우
            if e.errno in [errno.EACCES, errno.EPERM, 13, 32]:
                logger.debug(f"파일 잠금 감지 (IOError {e.errno}): {file_path}")
                return True
        except OSError as e:
            # Windows 특화: 다른 프로세스가 파일을 사용 중
            error_msg = str(e).lower()
            if e.errno in [errno.EACCES, errno.EPERM, 13, 32] or \
               'being used' in error_msg or \
               'locked' in error_msg or \
               'access denied' in error_msg or \
               'permission denied' in error_msg:
                logger.debug(f"파일 잠금 감지 (OSError): {file_path}")
                return True
        except Exception as e:
            # 예상치 못한 오류 - 안전하게 잠금으로 간주
            logger.debug(f"파일 체크 중 예외 발생 (안전하게 Skip): {file_path} - {e}")
            return True
        
        # 방법 2: Windows msvcrt를 사용한 추가 체크 (Python 3.8+)
        if sys.platform == 'win32':
            try:
                import msvcrt
                # 파일을 열어서 잠금 시도
                file_handle = open(file_path, 'rb')
                try:
                    # 파일 전체에 대한 잠금 시도 (non-blocking)
                    msvcrt.locking(file_handle.fileno(), msvcrt.LK_NBLCK, 1)
                    # 잠금 성공 - 즉시 해제
                    msvcrt.locking(file_handle.fileno(), msvcrt.LK_UNLCK, 1)
                    file_handle.close()
                except (IOError, OSError):
                    # 잠금 실패 - 다른 프로세스가 사용 중
                    file_handle.close()
                    logger.debug(f"파일 잠금 감지 (msvcrt): {file_path}")
                    return True
            except Exception:
                # msvcrt 체크 실패 - 기본값(안전) 사용
                pass
        
        # 모든 체크 통과 - 파일이 잠겨있지 않음
        return False
    
    
    def _extract_text(self, file_path: str) -> Optional[str]:
        """
        파일에서 텍스트 추출
        
        Args:
            file_path: 파일 경로
        
        Returns:
            추출된 텍스트 또는 None
        """
        ext = Path(file_path).suffix.lower()
        
        try:
            # 텍스트 파일
            if ext in self.SUPPORTED_TEXT_EXTENSIONS:
                return self._extract_text_file(file_path)
            
            # Word 문서
            elif ext == '.docx' and DOCX_AVAILABLE:
                return self._extract_docx(file_path)
            elif ext == '.doc' and WIN32COM_AVAILABLE:
                return self._extract_doc(file_path)
            
            # PowerPoint
            elif ext == '.pptx' and PPTX_AVAILABLE:
                return self._extract_pptx(file_path)
            elif ext == '.ppt' and WIN32COM_AVAILABLE:
                return self._extract_ppt(file_path)
            
            # Excel
            elif ext == '.xlsx' and XLSX_AVAILABLE:
                return self._extract_xlsx(file_path)
            elif ext == '.xls' and WIN32COM_AVAILABLE:
                return self._extract_xls(file_path)
            
            # CSV
            elif ext == '.csv':
                return self._extract_csv(file_path)
            
            # PDF
            elif ext == '.pdf' and PDF_AVAILABLE:
                return self._extract_pdf(file_path)
            
            # HWP
            elif ext == '.hwp':
                return self._extract_hwp(file_path)
        
        except Exception as e:
            logger.error(f"텍스트 추출 오류 [{file_path}]: {e}")
        
        return None
    
    def _extract_text_file(self, file_path: str) -> Optional[str]:
        """
        텍스트 파일 읽기 (인코딩 자동 감지)
        UTF-8 → CP949 → chardet 순서로 시도
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                logger.info(f"⛔ 텍스트 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일에서 텍스트 추출
            # 1차 시도: UTF-8
            try:
                with open(temp_file, 'r', encoding='utf-8') as f:
                    content = f.read()
                    return content[:100000]  # 최대 100KB
            except (UnicodeDecodeError, UnicodeError):
                pass
            
            # 2차 시도: CP949 (한글 Windows 기본 인코딩)
            try:
                with open(temp_file, 'r', encoding='cp949') as f:
                    content = f.read()
                    return content[:100000]
            except (UnicodeDecodeError, UnicodeError):
                pass
            
            # 3차 시도: chardet 자동 감지
            with open(temp_file, 'rb') as f:
                raw_data = f.read(1000000)  # 최대 1MB 읽기
                result = chardet.detect(raw_data)
                encoding = result['encoding']
                
                if encoding:
                    try:
                        content = raw_data.decode(encoding, errors='ignore')
                        return content[:100000]
                    except Exception:
                        pass
            
            # 최종: ignore 모드로 UTF-8 시도
            with open(temp_file, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                return content[:100000]
        
        except Exception as e:
            logger.debug(f"텍스트 파일 읽기 오류 [{file_path}]: {e}")
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_docx(self, file_path: str) -> Optional[str]:
        """
        Word 문서에서 텍스트 추출
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                logger.info(f"⛔ DOCX 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일에서 텍스트 추출
            doc = docx.Document(temp_file)
            text = '\n'.join([para.text for para in doc.paragraphs])
            
            logger.debug(f"✅ DOCX 파일 인덱싱 완료 (임시 복사본): {os.path.basename(file_path)}")
            
            return text[:100000]
            
        except Exception as e:
            error_msg = str(e).lower()
            filename = os.path.basename(file_path)
            
            # 암호화/손상된 파일 = 영구 Skip (재시도 X)
            if any(keyword in error_msg for keyword in ['password', 'encrypted', 'protected', 'corrupt', 'invalid', 'bad zipfile']):
                logger.info(f"⛔ DOCX 인덱싱 불가 (암호화/손상) - 영구 Skip: {filename}")
                self._log_skip(file_path, f"인덱싱 불가능 (암호화/손상): {str(e)[:80]}")
            # 접근 불가 파일 = 나중에 재시도
            elif 'being used' in error_msg or 'locked' in error_msg or 'permission denied' in error_msg:
                logger.info(f"⛔ DOCX 파일 접근 불가 - 나중에 재시도: {filename}")
                self._log_skip(file_path, "파일 접근 불가 - 재시도 예정")
                self._add_to_retry_queue(file_path, "파일 접근 불가")
            else:
                logger.debug(f"DOCX 추출 오류 [{filename}]: {e}")
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_pptx(self, file_path: str) -> Optional[str]:
        """
        PowerPoint에서 텍스트 추출
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                logger.info(f"⛔ PPTX 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일에서 텍스트 추출
            prs = Presentation(temp_file)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            
            logger.debug(f"✅ PPTX 파일 인덱싱 완료 (임시 복사본): {os.path.basename(file_path)}")
            
            return '\n'.join(text_parts)[:100000]
            
        except Exception as e:
            error_msg = str(e).lower()
            filename = os.path.basename(file_path)
            
            # 암호화/손상된 PPTX = 영구 Skip
            if any(keyword in error_msg for keyword in ['password', 'encrypted', 'protected', 'corrupt', 'invalid', 'bad zipfile']):
                logger.info(f"⛔ PPTX 인덱싱 불가 (암호화/손상) - 영구 Skip: {filename}")
                self._log_skip(file_path, f"인덱싱 불가능 (암호화/손상): {str(e)[:80]}")
            # 접근 불가 = 재시도
            elif 'being used' in error_msg or 'locked' in error_msg or 'permission denied' in error_msg:
                logger.info(f"⛔ PPTX 파일 접근 불가 - 나중에 재시도: {filename}")
                self._log_skip(file_path, "파일 접근 불가 - 재시도 예정")
                self._add_to_retry_queue(file_path, "파일 접근 불가")
            else:
                logger.debug(f"PPTX 추출 오류 [{filename}]: {e}")
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_doc(self, file_path: str) -> Optional[str]:
        """
        구버전 Word 문서(.doc)에서 텍스트 추출
        pywin32 COM 객체 사용 (Windows 전용)
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사 (사용자 파일 보호)
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                # 복사 실패 (파일이 잠겨있거나 접근 불가)
                logger.info(f"⛔ DOC 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일로 COM 작업 (원본 파일은 절대 건드리지 않음)
            pythoncom.CoInitialize()
            
            # DispatchEx로 완전히 새로운 Word 인스턴스 생성 (사용자 Word와 격리)
            word = win32com.client.DispatchEx("Word.Application")
            word.Visible = False
            word.DisplayAlerts = False
            
            # 임시 파일 열기 (원본 파일 X)
            doc = word.Documents.Open(
                temp_file,  # 임시 파일 사용!
                ReadOnly=True,
                ConfirmConversions=False,
                AddToRecentFiles=False
            )
            text = doc.Content.Text
            doc.Close(False)
            word.Quit()
            
            pythoncom.CoUninitialize()
            
            logger.info(f"✅ DOC 파일 인덱싱 완료 (임시 복사본 사용): {os.path.basename(file_path)}")
            
            return text[:100000]
            
        except Exception as e:
            error_msg = str(e).lower()
            if 'being used' in error_msg or 'locked' in error_msg or 'access denied' in error_msg:
                logger.info(f"⛔ DOC 파일 접근 불가 - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일 접근 불가")
            else:
                logger.debug(f"DOC 추출 오류 [{file_path}]: {e}")
            
            try:
                word.Quit()
            except:
                pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_ppt(self, file_path: str) -> Optional[str]:
        """
        구버전 PowerPoint(.ppt)에서 텍스트 추출
        pywin32 COM 객체 사용 (Windows 전용)
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                logger.info(f"⛔ PPT 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일로 COM 작업
            pythoncom.CoInitialize()
            
            # DispatchEx로 완전히 새로운 PowerPoint 인스턴스 생성 (사용자 PowerPoint와 격리)
            ppt = win32com.client.DispatchEx("PowerPoint.Application")
            ppt.Visible = False
            ppt.DisplayAlerts = False
            
            # 임시 파일 열기
            presentation = ppt.Presentations.Open(temp_file, ReadOnly=True, WithWindow=False)
            text_parts = []
            
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    if hasattr(shape, "TextFrame"):
                        if hasattr(shape.TextFrame, "TextRange"):
                            text_parts.append(shape.TextFrame.TextRange.Text)
            
            presentation.Close()
            ppt.Quit()
            
            pythoncom.CoUninitialize()
            
            logger.info(f"✅ PPT 파일 인덱싱 완료 (임시 복사본 사용): {os.path.basename(file_path)}")
            
            return '\n'.join(text_parts)[:100000]
            
        except Exception as e:
            error_msg = str(e).lower()
            if 'being used' in error_msg or 'locked' in error_msg or 'access denied' in error_msg:
                logger.info(f"⛔ PPT 파일 접근 불가 - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일 접근 불가")
            else:
                logger.debug(f"PPT 추출 오류 [{file_path}]: {e}")
            
            try:
                ppt.Quit()
            except:
                pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_xlsx(self, file_path: str) -> Optional[str]:
        """
        Excel 문서(.xlsx)에서 텍스트 추출
        openpyxl 사용, data_only=True로 수식 제외 값만 추출
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                logger.info(f"⛔ XLSX 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일에서 텍스트 추출
            workbook = openpyxl.load_workbook(temp_file, data_only=True, read_only=True)
            text_parts = []
            
            # 모든 시트 순회
            for sheet_name in workbook.sheetnames:
                # 시트 이름 추가 (검색 가능하도록)
                text_parts.append(f"\n[시트: {sheet_name}]\n")
                
                sheet = workbook[sheet_name]
                
                # 모든 행 순회
                for row in sheet.iter_rows(values_only=True):
                    for cell_value in row:
                        if cell_value is not None:
                            text_parts.append(str(cell_value))
            
            workbook.close()
            
            logger.debug(f"✅ XLSX 파일 인덱싱 완료 (임시 복사본): {os.path.basename(file_path)}")
            
            return ' '.join(text_parts)[:100000]
            
        except Exception as e:
            error_msg = str(e).lower()
            filename = os.path.basename(file_path)
            
            # 암호화/손상된 Excel = 영구 Skip
            if any(keyword in error_msg for keyword in ['password', 'encrypted', 'protected', 'corrupt', 'invalid', 'bad zipfile']):
                logger.info(f"⛔ XLSX 인덱싱 불가 (암호화/손상) - 영구 Skip: {filename}")
                self._log_skip(file_path, f"인덱싱 불가능 (암호화/손상): {str(e)[:80]}")
            # 접근 불가 = 재시도
            elif 'being used' in error_msg or 'locked' in error_msg or 'permission denied' in error_msg:
                logger.info(f"⛔ XLSX 파일 접근 불가 - 나중에 재시도: {filename}")
                self._log_skip(file_path, "파일 접근 불가 - 재시도 예정")
                self._add_to_retry_queue(file_path, "파일 접근 불가")
            else:
                logger.debug(f"XLSX 추출 오류 [{filename}]: {e}")
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_xls(self, file_path: str) -> Optional[str]:
        """
        레거시 Excel(.xls)에서 텍스트 추출
        pywin32 COM 객체 사용 (Windows 전용)
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                logger.info(f"⛔ XLS 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일로 COM 작업
            pythoncom.CoInitialize()
            
            # DispatchEx로 완전히 새로운 Excel 인스턴스 생성 (사용자 Excel과 격리)
            excel = win32com.client.DispatchEx("Excel.Application")
            excel.Visible = False
            excel.DisplayAlerts = False
            
            # 임시 파일 열기
            workbook = excel.Workbooks.Open(temp_file, ReadOnly=True)
            text_parts = []
            
            # 모든 시트 순회
            for sheet in workbook.Sheets:
                # 시트 이름 추가 (검색 가능하도록)
                sheet_name = sheet.Name
                text_parts.append(f"\n[시트: {sheet_name}]\n")
                
                used_range = sheet.UsedRange
                for row in used_range.Rows:
                    for cell in row.Cells:
                        if cell.Value is not None:
                            text_parts.append(str(cell.Value))
            
            workbook.Close(False)
            excel.Quit()
            
            pythoncom.CoUninitialize()
            
            logger.info(f"✅ XLS 파일 인덱싱 완료 (임시 복사본 사용): {os.path.basename(file_path)}")
            
            return ' '.join(text_parts)[:100000]
            
        except Exception as e:
            error_msg = str(e).lower()
            if 'being used' in error_msg or 'locked' in error_msg or 'access denied' in error_msg:
                logger.info(f"⛔ XLS 파일 접근 불가 - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일 접근 불가")
            else:
                logger.debug(f"XLS 추출 오류 [{file_path}]: {e}")
            
            try:
                excel.Quit()
            except:
                pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_csv(self, file_path: str) -> Optional[str]:
        """
        CSV 파일에서 텍스트 추출
        Python 기본 csv 모듈 사용
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        import csv
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                logger.info(f"⛔ CSV 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일에서 텍스트 추출
            text_parts = []
            
            # UTF-8, CP949(한글), Latin-1 순서로 시도
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
            content_read = False
            
            for encoding in encodings:
                try:
                    with open(temp_file, 'r', encoding=encoding, errors='ignore') as f:
                        csv_reader = csv.reader(f)
                        
                        for row in csv_reader:
                            # 각 행의 셀들을 탭으로 구분하여 추가
                            row_text = '\t'.join(str(cell) for cell in row if cell)
                            if row_text.strip():
                                text_parts.append(row_text)
                    
                    content_read = True
                    logger.debug(f"✅ CSV 파일 인덱싱 완료 (임시 복사본, 인코딩: {encoding}): {os.path.basename(file_path)}")
                    break
                    
                except (UnicodeDecodeError, LookupError):
                    continue
            
            if not content_read:
                logger.info(f"⚠️ CSV 파일 인코딩 처리 실패: {os.path.basename(file_path)}")
                return None
            
            return '\n'.join(text_parts)[:100000]
            
        except Exception as e:
            error_msg = str(e).lower()
            filename = os.path.basename(file_path)
            
            # 접근 불가 = 재시도
            if 'being used' in error_msg or 'locked' in error_msg or 'permission denied' in error_msg:
                logger.info(f"⛔ CSV 파일 접근 불가 - 나중에 재시도: {filename}")
                self._log_skip(file_path, "파일 접근 불가 - 재시도 예정")
                self._add_to_retry_queue(file_path, "파일 접근 불가")
            else:
                logger.debug(f"CSV 추출 오류 [{filename}]: {e}")
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_pdf(self, file_path: str) -> Optional[str]:
        """
        PDF에서 텍스트 추출
        PyMuPDF (fitz) 사용 - 속도가 월등히 빠름
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        """
        temp_file = None
        
        try:
            # 1단계: 원본 파일을 임시 폴더에 복사
            temp_file = self._copy_to_temp(file_path)
            
            if not temp_file:
                logger.info(f"⛔ PDF 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                return None
            
            # 2단계: 임시 파일에서 텍스트 추출
            doc = fitz.open(temp_file)
            text_parts = []
            
            # 최대 100페이지까지만
            for page_num in range(min(len(doc), 100)):
                page = doc[page_num]
                text_parts.append(page.get_text())
            
            doc.close()
            
            logger.debug(f"✅ PDF 파일 인덱싱 완료 (임시 복사본): {os.path.basename(file_path)}")
            
            return '\n'.join(text_parts)[:100000]
            
        except Exception as e:
            error_msg = str(e).lower()
            filename = os.path.basename(file_path)
            
            # 암호화/손상된 PDF = 영구 Skip
            if any(keyword in error_msg for keyword in ['password', 'encrypted', 'protected', 'corrupt', 'invalid', 'damaged']):
                logger.info(f"⛔ PDF 인덱싱 불가 (암호화/손상) - 영구 Skip: {filename}")
                self._log_skip(file_path, f"인덱싱 불가능 (암호화/손상): {str(e)[:80]}")
            # 접근 불가 = 재시도
            elif 'being used' in error_msg or 'locked' in error_msg or 'permission denied' in error_msg:
                logger.info(f"⛔ PDF 파일 접근 불가 - 나중에 재시도: {filename}")
                self._log_skip(file_path, "파일 접근 불가 - 재시도 예정")
                self._add_to_retry_queue(file_path, "파일 접근 불가")
            else:
                logger.debug(f"PDF 추출 오류 [{filename}]: {e}")
            return None
            
        finally:
            # 3단계: 임시 파일 정리
            if temp_file:
                self._cleanup_temp(temp_file)
    
    def _extract_hwp(self, file_path: str) -> Optional[str]:
        """
        HWP 파일에서 텍스트 추출
        1차: pywin32 COM 객체 시도 (임시 파일 사용 + 타임아웃)
        2차: olefile 라이브러리 시도
        
        🛡️ 안전 모드: 원본 파일을 건드리지 않고 임시 복사본으로 인덱싱합니다!
        ⏱️ 타임아웃: 30초 이상 걸리면 자동 Skip
        """
        temp_file = None
        hwp_timeout = 30  # HWP 파일 처리 타임아웃: 30초
        
        # 1차 시도: COM 객체 (가장 정확)
        if WIN32COM_AVAILABLE:
            try:
                # 임시 파일 복사
                temp_file = self._copy_to_temp(file_path)
                
                if not temp_file:
                    logger.info(f"⛔ HWP 파일 복사 실패 (사용 중) - Skip: {os.path.basename(file_path)}")
                    self._log_skip(file_path, "파일이 사용 중이거나 접근 불가")
                    return None
                
                # COM 작업을 스레드에서 실행하여 타임아웃 적용
                result_container = [None]
                error_container = [None]
                
                def hwp_extract_thread():
                    try:
                        pythoncom.CoInitialize()
                        
                        # DispatchEx로 완전히 새로운 한글 인스턴스 생성 (사용자 한글과 격리)
                        hwp = win32com.client.DispatchEx("HWPFrame.HwpObject")
                        hwp.RegisterModule("FilePathCheckDLL", "SecurityModule")
                        hwp.Open(temp_file)  # 임시 파일 사용!
                        
                        hwp.InitScan()
                        text_parts = []
                        
                        while True:
                            text = hwp.GetText()
                            if not text:
                                break
                            text_parts.append(text)
                        
                        hwp.ReleaseScan()
                        hwp.Quit()
                        
                        pythoncom.CoUninitialize()
                        
                        result_container[0] = ''.join(text_parts)[:100000]
                        
                    except Exception as e:
                        error_container[0] = e
                        try:
                            hwp.Quit()
                        except:
                            pass
                        try:
                            pythoncom.CoUninitialize()
                        except:
                            pass
                
                # 스레드 시작 및 타임아웃 대기
                import threading
                thread = threading.Thread(target=hwp_extract_thread, daemon=True)
                thread.start()
                thread.join(timeout=hwp_timeout)
                
                if thread.is_alive():
                    # 타임아웃 발생
                    logger.warning(f"⏰ HWP 파일 처리 타임아웃 ({hwp_timeout}초) - Skip: {os.path.basename(file_path)}")
                    self._log_skip(file_path, f"HWP 파일 처리 타임아웃 ({hwp_timeout}초)")
                    if temp_file:
                        self._cleanup_temp(temp_file)
                    return None
                
                # 오류 확인
                if error_container[0]:
                    raise error_container[0]
                
                # 성공
                if result_container[0]:
                    logger.info(f"✅ HWP 파일 인덱싱 완료 (임시 복사본 사용): {os.path.basename(file_path)}")
                    if temp_file:
                        self._cleanup_temp(temp_file)
                    return result_container[0]
                
            except Exception as e:
                logger.debug(f"HWP COM 추출 오류 [{file_path}]: {e}")
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
                # 임시 파일 정리
                if temp_file:
                    self._cleanup_temp(temp_file)
        
        # 2차 시도: olefile (제한적) - 임시 파일 사용
        if OLEFILE_AVAILABLE:
            try:
                if not temp_file:
                    temp_file = self._copy_to_temp(file_path)
                
                if not temp_file:
                    return None
                
                ole = olefile.OleFileIO(temp_file)
                if ole.exists('PrvText'):
                    stream = ole.openstream('PrvText')
                    data = stream.read()
                    # HWP 텍스트는 UTF-16LE 인코딩
                    text = data.decode('utf-16le', errors='ignore')
                    ole.close()
                    
                    logger.debug(f"✅ HWP 파일 인덱싱 완료 (olefile, 임시 복사본): {os.path.basename(file_path)}")
                    
                    # 임시 파일 정리
                    if temp_file:
                        self._cleanup_temp(temp_file)
                    
                    return text[:100000]
                ole.close()
            except Exception as e:
                logger.debug(f"HWP olefile 추출 오류 [{file_path}]: {e}")
            finally:
                # 임시 파일 정리
                if temp_file:
                    self._cleanup_temp(temp_file)
        
        logger.debug(f"HWP 파일 추출 실패 [{file_path}]: 지원 라이브러리 없음")
        
        # 마지막 정리
        if temp_file:
            self._cleanup_temp(temp_file)
        
        return None
    
    def get_stats(self) -> dict:
        """인덱싱 통계 반환"""
        return self.stats.copy()
    
    def start_retry_worker(self):
        """
        재시도 워커 시작 (백그라운드 스레드)
        
        5-10분마다 Skip된 파일을 재시도하여 인덱싱
        """
        if self.retry_thread and self.retry_thread.is_alive():
            logger.warning("재시도 워커가 이미 실행 중입니다.")
            return
        
        self.retry_stop_flag.clear()
        self.retry_thread = threading.Thread(
            target=self._retry_worker,
            name="RetryWorker",
            daemon=True
        )
        self.retry_thread.start()
        logger.info(f"재시도 워커 시작됨 (간격: {self.retry_interval}초)")
    
    def _retry_worker(self):
        """
        재시도 워커 스레드
        
        주기적으로 Skip된 파일을 재시도하여 인덱싱
        """
        logger.info("재시도 워커 동작 시작")
        
        while not self.retry_stop_flag.is_set():
            # 대기 (5분 = 300초, 인터럽트 가능하도록 1초씩 체크)
            for _ in range(self.retry_interval):
                if self.retry_stop_flag.is_set():
                    break
                time.sleep(1)
            
            if self.retry_stop_flag.is_set():
                break
            
            # Skip된 파일 재시도
            with self.skipped_files_lock:
                if not self.skipped_files:
                    logger.info("재시도할 파일이 없습니다. 워커 종료.")
                    break
                
                files_to_retry = list(self.skipped_files.keys())
            
            logger.info(f"Skip된 파일 재시도 시작: {len(files_to_retry)}개")
            
            retry_success = 0
            retry_failed = 0
            
            for file_path in files_to_retry:
                if self.retry_stop_flag.is_set():
                    break
                
                # 사용자 활동 체크 (재시도 워커에도 적용)
                if self.activity_monitor and self.enable_activity_monitor:
                    if self.activity_monitor.is_user_active():
                        # 사용자 활동 감지 - 즉시 대기
                        logger.debug(f"⏸️ 재시도 워커: 사용자 작업 중 - 대기...")
                        # 유휴 상태가 될 때까지 대기 (더 짧은 체크 간격)
                        if not self.activity_monitor.wait_until_idle(check_interval=0.1, stop_flag=self.retry_stop_flag):
                            # 중지 요청됨
                            break
                        logger.debug("▶️ 재시도 워커: 사용자 활동 없음 - 재개")
                
                try:
                    # 파일이 존재하는지 확인
                    if not os.path.exists(file_path):
                        with self.skipped_files_lock:
                            if file_path in self.skipped_files:
                                del self.skipped_files[file_path]
                        logger.debug(f"파일 삭제됨, 재시도 목록에서 제거: {file_path}")
                        continue
                    
                    # 파일 크기 재확인
                    try:
                        file_size = os.path.getsize(file_path)
                        if file_size > MAX_FILE_SIZE:
                            with self.skipped_files_lock:
                                if file_path in self.skipped_files:
                                    del self.skipped_files[file_path]
                            logger.debug(f"파일 크기 초과, 재시도 중단: {file_path}")
                            continue
                    except Exception:
                        pass
                    
                    # 텍스트 추출 재시도
                    content = self._extract_text_safe(file_path)
                    
                    if content:
                        # 성공! DB에 저장
                        current_mtime = os.path.getmtime(file_path)
                        token_count = self._count_tokens(content)
                        
                        # 이미 DB에 있는지 확인
                        indexed_mtime = self.db.get_file_mtime(file_path)
                        
                        try:
                            if indexed_mtime is not None:
                                # 업데이트
                                self.db.update_file(file_path, content, current_mtime)
                            else:
                                # 새로 삽입
                                self.db.insert_file(file_path, content, current_mtime)
                            
                            # 재시도 목록에서 제거
                            with self.skipped_files_lock:
                                if file_path in self.skipped_files:
                                    retry_info = self.skipped_files[file_path]
                                    del self.skipped_files[file_path]
                                    logger.info(f"재시도 성공 [{file_path}] - 이전 사유: {retry_info['reason']}")
                            
                            retry_success += 1
                            
                            # UI 로그 콜백 및 메모리에 로그 추가 - DB 저장 완료 상태
                            filename = os.path.basename(file_path)
                            db_status = "✓ DB 저장 완료 (재시도)"
                            detail = f'{len(content):,}자 / {token_count:,}토큰 | {db_status}'
                            
                            # 통합 로그에 기록
                            self._write_indexing_log('Retry Success', file_path, detail)
                            
                            # Indexed.txt에 기록 (재시도 성공도 인덱싱 성공)
                            self._write_indexed_file(file_path, len(content), token_count, content)
                            
                            self._add_log_to_memory('Retry Success', file_path, detail)
                            
                            if self.log_callback:
                                self.log_callback('Retry Success', filename, detail)
                        
                        except Exception as e:
                            logger.error(f"재시도 DB 저장 오류 [{file_path}]: {e}")
                            retry_failed += 1
                    
                    else:
                        # 여전히 실패 - 무제한 재시도 (요구사항: 사용자가 사용중이면 절대 프로그램을 닫지 않도록 함)
                        with self.skipped_files_lock:
                            if file_path in self.skipped_files:
                                self.skipped_files[file_path]['retry_count'] += 1
                                retry_count = self.skipped_files[file_path]['retry_count']
                                logger.debug(f"재시도 실패 (재시도 횟수: {retry_count}회): {file_path}")
                        
                        retry_failed += 1
                
                except Exception as e:
                    logger.error(f"재시도 중 오류 [{file_path}]: {e}")
                    retry_failed += 1
            
            logger.info(f"재시도 완료: 성공 {retry_success}개, 실패 {retry_failed}개")
        
        logger.info("재시도 워커 종료")
    
    def get_skipped_files_count(self) -> int:
        """현재 재시도 대기 중인 파일 수 반환"""
        with self.skipped_files_lock:
            return len(self.skipped_files)


# 테스트 코드
if __name__ == "__main__":
    # 데이터베이스 초기화
    db = DatabaseManager("test_index.db")
    
    # 인덱서 생성
    indexer = FileIndexer(db)
    
    # 진행 상황 콜백
    def progress(current, total, path):
        print(f"[{current}/{total}] {path}")
    
    # 인덱싱 시작
    test_path = os.path.expanduser("~\\Desktop")
    print(f"테스트 인덱싱: {test_path}")
    indexer.start_indexing([test_path], progress)
    
    # 완료 대기
    while indexer.is_running:
        time.sleep(1)
    
    # 통계 출력
    stats = indexer.get_stats()
    print(f"\n=== 인덱싱 통계 ===")
    print(f"총 파일: {stats['total_files']}")
    print(f"인덱싱됨: {stats['indexed_files']}")
    print(f"스킵됨: {stats['skipped_files']}")
    print(f"오류: {stats['error_files']}")
    
    # 검색 테스트
    results = db.search("test")
    print(f"\n검색 결과: {len(results)}개")
    
    # 정리
    db.close()
    os.remove("test_index.db")
    print("\n테스트 완료")

